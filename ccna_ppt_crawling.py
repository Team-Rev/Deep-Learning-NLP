from collections import defaultdict
from pptx import Presentation
import os
from sklearn.feature_extraction.text import CountVectorizer, TfidfVectorizer
import re
from nltk.tokenize import word_tokenize
import matplotlib.pyplot as plt
from nltk.corpus import stopwords
import pandas as pd
from collections import Counter
from gensim.models import Word2Vec
from wordcloud import WordCloud


text = []
stop_words = stopwords.words('english')
stop_words.remove('a')

root_dir = os.listdir('CCNA')
for i in range(len(root_dir)):
    sub_dir = os.listdir('CCNA/'+f'{root_dir[i]}')
    for j in range(len(sub_dir)):
        prs = Presentation(f'CCNA/{root_dir[i]}/{sub_dir[j]}')
        slide_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                has_table = shape.has_table
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.text = re.sub(r"^\s+", "", paragraph.text)
                        slide_text.append(paragraph.text)
                elif shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            cell.text = re.sub(r"^\s+", "", cell.text)
                            slide_text.append(cell.text)
        text.append(str(slide_text))

for token in text:
    if len(token) < 1:
        text.remove(token)

data = pd.Series(text)


data = data.str.replace("[^a-zA-Z0-9]", " ", regex=True)
data = data.apply(lambda x: ' '.join([w for w in x.split() if w not in stop_words]))

count = 1
for chapter in data:
    text = str(chapter)
    wordcloud = WordCloud()
    wordcloud = wordcloud.generate(text)
    plt.figure(figsize=(22, 22))  # 이미지 사이즈 지정
    plt.imshow(wordcloud, interpolation='lanczos')  # 이미지의 부드럽기 정도
    plt.axis('off')  # x y 축 숫자 제거
    plt.savefig(f'{count}_wordcloud.png')
    count += 1


# 토큰화
# tokenized = [word_tokenize(str(token).lower()) for token in data]






# BoW
# vector = CountVectorizer()
# print(vector.fit_transform(tokenized).toarray())
#
# print(vector.vocabulary_)

# TF-IDF
# tf = TfidfVectorizer()
# tf.fit(tokenized)
# tfidf = tf.transform(tokenized).toarray()
# word2id = defaultdict(lambda :0)
# for idx, feature in enumerate(tf.get_feature_names()):
#     word2id[feature] = idx
# for i, sent in enumerate(tokenized):
#     print('==== 문장[%d] ====' % i)
#     print([(token, tfidf[i, word2id[token]]) for token in sent.split()])



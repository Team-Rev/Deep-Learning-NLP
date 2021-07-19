from nltk import word_tokenize

import csv

from pptx import Presentation
import os
from nltk.corpus import stopwords
import re
import pandas as pd
import load_phrase

custom_stop_words = []
f = open('stopword.txt', 'r')
category = open('category.txt', 'w')
while True:
    line = f.readline()
    line = line.strip()
    if not line: break
    custom_stop_words.append(line)
f.close()

text = []
stop_words = stopwords.words('english')
stop_words.extend(custom_stop_words)
load_phrase = load_phrase.load_phrase()
phrases = load_phrase.get_phrases()

root_dir = os.listdir('CCNA/Category_classfication')

w = open('test.csv', 'w', newline='')
wr = csv.writer(w)

for i in range(len(root_dir)):
    sub_dir = os.listdir('CCNA/Category_classfication'+f'/{root_dir[i]}')
    category.write(root_dir[i]+'\n')
    for j in range(len(sub_dir)):
        prs = Presentation(f'CCNA/Category_classfication/{root_dir[i]}/{sub_dir[j]}')
        slide_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.text = re.sub(r"^\s+", "", paragraph.text.lower())
                        paragraph.text = paragraph.text.replace('\\xa0', ' ')
                        paragraph.text = paragraph.text.replace('\\x0', ' ')
                        paragraph.text = re.sub(r"[^a-zA-Z0-9-]", " ", paragraph.text)
                        slide_text.append(paragraph.text)
                elif shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            cell.text = re.sub(r"^\s+", "", cell.text.lower())
                            cell.text = cell.text.replace('\\xa0', ' ')
                            cell.text = cell.text.replace('\\x0', ' ')
                            cell.text = re.sub(r"[^a-zA-Z0-9-\\s]", " ", cell.text)
                            slide_text.append(cell.text)
        chapter_text = ' '.join(slide_text)
        for phrase in phrases:
            chapter_text = chapter_text.replace(phrase, re.sub(' +', '-', phrase))
        text.append(' '.join([w for w in chapter_text.split() if w not in stop_words and len(w) > 2]))


category = pd.DataFrame()

for i, category in enumerate(text):
    tokenized = word_tokenize(category)
    tokenized = list(dict.fromkeys(tokenized))





